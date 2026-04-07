"""
DDVB Naming Presentation PPTX Generator — Template

Usage: Adapt this script per project. Replace placeholders with actual data.
Requires: pip install python-pptx

DDVB Brand System:
- Background: #1A1A1A (dark)
- Accent: #FFC000 (orange/gold)
- Text primary: #FFFFFF
- Text secondary: #888888
- Fonts: Atyp Display (headings), Atyp Text (body) — fallback to Calibri
- Corner markers: D/D/V/B in corners of each slide
"""

import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Constants ──
BG_COLOR = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT = RGBColor(0xFF, 0xC0, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_SURFACE = RGBColor(0x22, 0x22, 0x22)
HEADING_FONT = "Atyp Display"
BODY_FONT = "Atyp Text"
FALLBACK_FONT = "Calibri"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def set_slide_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_corner_markers(slide):
    """Add D/D/V/B corner markers — DDVB brand element."""
    corners = [
        ("D", Inches(0.3), Inches(0.3)),
        ("D", Inches(12.5), Inches(0.3)),
        ("V", Inches(0.3), Inches(6.8)),
        ("B", Inches(12.5), Inches(6.8)),
    ]
    for letter, left, top in corners:
        txBox = slide.shapes.add_textbox(left, top, Inches(0.5), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = letter
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        try:
            p.font.name = BODY_FONT
        except:
            p.font.name = FALLBACK_FONT


def add_text(slide, text, left, top, width, height,
             font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    try:
        p.font.name = font_name or HEADING_FONT
    except:
        p.font.name = FALLBACK_FONT
    return txBox


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


# ── Slide Builders ──

def slide_title(prs, client_name, date):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    add_corner_markers(slide)

    add_text(slide, "NAMING", Inches(1.5), Inches(2.0), Inches(10), Inches(0.6),
             font_size=14, color=ACCENT, bold=True)
    add_accent_line(slide, Inches(1.5), Inches(2.7), Inches(2))
    add_text(slide, f"Нейминг для {client_name}", Inches(1.5), Inches(2.9), Inches(10), Inches(1.2),
             font_size=40, color=WHITE, bold=False)
    add_text(slide, date, Inches(1.5), Inches(4.3), Inches(10), Inches(0.4),
             font_size=14, color=GRAY, font_name=BODY_FONT)


def slide_process(prs):
    """Slide 2: Process overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_corner_markers(slide)

    add_text(slide, "ПРОЦЕСС", Inches(1.5), Inches(1.0), Inches(10), Inches(0.5),
             font_size=12, color=ACCENT, bold=True)
    add_text(slide, "Как мы пришли к этим кандидатам", Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
             font_size=28, color=WHITE)

    steps = [
        ("01", "Исследование", "Конкурентный анализ, лингвистический аудит"),
        ("02", "Бриф для копирайтера", "Семантическое поле, стратегии, критерии"),
        ("03", "Генерация", "AI-кандидаты + список от копирайтера"),
        ("04", "Оценка", "10 параметров, взвешенный скоринг"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        left = Inches(1.5 + i * 2.8)
        top = Inches(3.0)
        # Number
        add_text(slide, num, left, top, Inches(1), Inches(0.5),
                 font_size=32, color=ACCENT, bold=False)
        # Title
        add_text(slide, title, left, Inches(3.6), Inches(2.4), Inches(0.4),
                 font_size=14, color=WHITE, bold=True, font_name=BODY_FONT)
        # Desc
        add_text(slide, desc, left, Inches(4.1), Inches(2.4), Inches(0.8),
                 font_size=11, color=GRAY, font_name=BODY_FONT)


def slide_candidate(prs, name, strategy, etymology, strengths, availability):
    """Slide N: Individual candidate"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_corner_markers(slide)

    # Name — large, centered
    add_text(slide, name, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
             font_size=56, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5.5), Inches(3.2), Inches(2.3))

    # Strategy type
    add_text(slide, strategy.upper(), Inches(1.5), Inches(3.5), Inches(10), Inches(0.4),
             font_size=11, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    # Etymology
    add_text(slide, etymology, Inches(2.5), Inches(4.1), Inches(8), Inches(0.8),
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER, font_name=BODY_FONT)

    # Strengths (left column)
    add_text(slide, "ПРЕИМУЩЕСТВА", Inches(1.5), Inches(5.2), Inches(4), Inches(0.3),
             font_size=10, color=ACCENT, bold=True, font_name=BODY_FONT)
    strengths_text = "\n".join(f"• {s}" for s in strengths)
    add_text(slide, strengths_text, Inches(1.5), Inches(5.6), Inches(4), Inches(1.5),
             font_size=12, color=WHITE, font_name=BODY_FONT)

    # Availability (right column)
    add_text(slide, "ДОСТУПНОСТЬ", Inches(7.5), Inches(5.2), Inches(4), Inches(0.3),
             font_size=10, color=ACCENT, bold=True, font_name=BODY_FONT)
    avail_text = "\n".join(f"• {a}" for a in availability)
    add_text(slide, avail_text, Inches(7.5), Inches(5.6), Inches(4), Inches(1.5),
             font_size=12, color=WHITE, font_name=BODY_FONT)


def slide_comparison(prs, candidates_data):
    """Slide: Side-by-side comparison matrix"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_corner_markers(slide)

    add_text(slide, "СРАВНЕНИЕ", Inches(1.5), Inches(0.8), Inches(10), Inches(0.5),
             font_size=12, color=ACCENT, bold=True)
    add_text(slide, "Финальные кандидаты", Inches(1.5), Inches(1.3), Inches(10), Inches(0.6),
             font_size=24, color=WHITE)

    # Table would go here — adapt per project
    # candidates_data: list of {name, score, strengths_summary}


def slide_recommendation(prs, recommended_name, rationale):
    """Slide: Final recommendation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_corner_markers(slide)

    add_text(slide, "РЕКОМЕНДАЦИЯ", Inches(1.5), Inches(1.5), Inches(10), Inches(0.5),
             font_size=12, color=ACCENT, bold=True)

    add_text(slide, recommended_name, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
             font_size=64, color=ACCENT, alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5), Inches(4.3), Inches(3.3))

    add_text(slide, rationale, Inches(2), Inches(4.8), Inches(9), Inches(1.5),
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER, font_name=BODY_FONT)


def slide_next_steps(prs, steps):
    """Slide: Next steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_corner_markers(slide)

    add_text(slide, "СЛЕДУЮЩИЕ ШАГИ", Inches(1.5), Inches(1.0), Inches(10), Inches(0.5),
             font_size=12, color=ACCENT, bold=True)
    add_text(slide, "Что дальше", Inches(1.5), Inches(1.5), Inches(10), Inches(0.6),
             font_size=28, color=WHITE)

    for i, step in enumerate(steps):
        num = f"0{i+1}"
        add_text(slide, num, Inches(1.5), Inches(2.8 + i * 1.0), Inches(0.8), Inches(0.4),
                 font_size=24, color=ACCENT)
        add_text(slide, step, Inches(2.5), Inches(2.85 + i * 1.0), Inches(9), Inches(0.4),
                 font_size=14, color=WHITE, font_name=BODY_FONT)


# ── Example Usage ──
if __name__ == "__main__":
    prs = create_presentation()

    slide_title(prs, "Client Name", "Апрель 2026")
    slide_process(prs)
    slide_candidate(prs,
        name="MERIDIAN",
        strategy="Metaphorical",
        etymology="Линия, соединяющая полюса — символ связи, масштаба, точности навигации",
        strengths=["Уникальная семантическая территория", "Работает на 3 языках", "Сильный визуальный потенциал"],
        availability=[".com — доступен", "Telegram — доступен", "Роспатент — чисто в классе 39"]
    )
    slide_recommendation(prs, "MERIDIAN", "Метафора точности и глобального масштаба. Свободная территория. Работает на всех целевых языках.")
    slide_next_steps(prs, [
        "Регистрация товарного знака (ФИПС, классы 35/39)",
        "Регистрация домена meridian.kz + meridian-logistics.com",
        "Разработка визуальной идентичности",
        "Запуск коммуникационной платформы",
    ])

    prs.save("EXAMPLE_NAMING_PRESENTATION.pptx")
    print("Saved: EXAMPLE_NAMING_PRESENTATION.pptx")
