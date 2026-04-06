#!/usr/bin/env python3
"""Build the final BCS repositioning presentation (15 slides, DDVB brand).
Slides 4, 10, 13, 14 are native PPTX (not PNG). All slides have headers."""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

SKILL_DIR = Path.home() / ".claude/skills/ddvb-brand-guidelines/resources"
BG_DIR = SKILL_DIR / "backgrounds"
FIG_DIR = Path(__file__).resolve().parent.parent / "figures"
OUT_DIR = Path(__file__).resolve().parent.parent

GOLD = RGBColor(0xFD, 0xB7, 0x1C)
BLACK = RGBColor(0x0D, 0x0D, 0x0D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x44, 0x54, 0x6A)
LGRAY = RGBColor(0xBD, 0xBD, 0xBD)
VLGRAY = RGBColor(0xD9, 0xD9, 0xD9)
BLUE = RGBColor(0x15, 0x65, 0xC0)
VIOLET = RGBColor(0x6A, 0x1B, 0x9A)
AMBER = RGBColor(0xE6, 0x51, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)

FD = "Atyp Display"
FT = "Atyp Text"


def add_bg(slide, prs, bg_name):
    pic = slide.shapes.add_picture(str(BG_DIR / bg_name), 0, 0,
                                    width=prs.slide_width, height=prs.slide_height)
    sp = pic._element
    slide.shapes._spTree.remove(sp)
    slide.shapes._spTree.insert(2, sp)


def txt(slide, l, t, w, h, text, font=FT, sz=18, color=BLACK,
        bold=False, align=PP_ALIGN.LEFT, lsp=1.2):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.name = font; p.font.size = Pt(sz); p.font.color.rgb = color
    p.font.bold = bold; p.alignment = align; p.line_spacing = Pt(sz * lsp)
    return tf


def header(slide, num, title, color=BLACK):
    """Standard slide header: '0X  Title'."""
    txt(slide, 0.8, 0.35, 8.4, 0.5, f"{num}  {title}", FD, 26, color, bold=True)


def rect(slide, l, t, w, h, fill_color, alpha_pct=100):
    """Add a filled rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.05
    return shape


def img(slide, path, l, t, w=None, h=None, **kw):
    args = {"image_file": str(path), "left": Inches(l), "top": Inches(t)}
    if w: args["width"] = Inches(w)
    if h: args["height"] = Inches(h)
    return slide.shapes.add_picture(**args)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    B = prs.slide_layouts[6]

    # ── SLIDE 1: Title ──────────────────────────────────────────────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_2.png")
    txt(s, 1.5, 1.4, 7, 1.2, "БКС Мир Инвестиций", FD, 44, WHITE, True, PP_ALIGN.CENTER)
    txt(s, 1.5, 2.4, 7, 0.8, "Стратегия репозиционирования бренда", FD, 24, WHITE, align=PP_ALIGN.CENTER)
    txt(s, 1.5, 3.3, 7, 0.5, "Гипотеза 2: Активатор благосостояния", FT, 16, WHITE, align=PP_ALIGN.CENTER)
    txt(s, 1.5, 4.4, 7, 0.4, "DDVB  \u00B7  Март 2026", FT, 12, WHITE, align=PP_ALIGN.CENTER)

    # ── SLIDE 2: Contents ───────────────────────────────────────────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_7.png")
    txt(s, 0.8, 0.5, 8, 0.6, "Содержание", FD, 32, BLACK, True)
    for i, (title, desc, w) in enumerate([
        ("01  Аналитический блок", "SWOT, конкурентный ландшафт, позиция на рынке", "25%"),
        ("02  Стратегический блок", "Платформа бренда: миссия, ценности, архетип, Big Idea", "25%"),
        ("03  Коммуникационный блок", "Ключевая идея, 3 сообщения, 5 слоганов, макеты", "20%"),
        ("04  Визуальный блок", "Айдентика, цвет, типографика, графический приём", "30%"),
    ]):
        y = 1.5 + i * 0.9
        txt(s, 0.8, y, 6.5, 0.35, title, FD, 18, BLACK, True)
        txt(s, 0.8, y + 0.35, 6.5, 0.3, desc, FT, 12, DGRAY)
        txt(s, 8.2, y + 0.1, 1.2, 0.35, w, FD, 20, GOLD, True, PP_ALIGN.RIGHT)

    # ── SLIDE 3: Market Position ────────────────────────────────────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_8.png")
    header(s, "01", "Позиция на рынке")
    for i, (val, lab) in enumerate([
        ("7,5%", "доля активов\nна Мосбирже"), ("25,4%", "доля оборотов\n(ист. максимум)"),
        ("8,16%", "доля клиентов\nс активами 2М+"), ("30 лет", "на рынке\nинвестиций"),
    ]):
        x = 0.8 + i * 2.3
        txt(s, x, 1.3, 2, 0.5, val, FD, 36, GOLD, True)
        txt(s, x, 2.0, 2, 0.5, lab, FT, 11, DGRAY)
    txt(s, 0.8, 2.9, 8.4, 0.35, "Крупнейший частный брокер РФ  \u00B7  Рейтинг ruА-  \u00B7  1М+ клиентов  \u00B7  Монолайнер", FT, 11, DGRAY)
    txt(s, 0.8, 3.4, 8.4, 0.35, "5-е место по узнаваемости (21%) \u2014 разрыв между объективной силой и восприятием", FT, 12, RED, True)
    txt(s, 0.8, 3.9, 8.4, 0.7, "Топ-3 по знанию в сегменте 2М+ \u2014 целевая аудитория знает БКС значительно лучше, чем массовый рынок. Конверсия в рассмотрение выше среднего.", FT, 11, DGRAY)

    # ── SLIDE 4: SWOT (NATIVE PPTX) ────────────────────────────────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_7.png")
    header(s, "01", "SWOT-анализ")

    swot_data = [
        ("STRENGTHS", GREEN, 0.5, 1.05, [
            "S1  Крупнейший частный брокер РФ",
            "S2  30 лет экспертизы, рейтинг ruА-",
            "S3  Топ-3 по знанию в сегменте 2М+",
            "S4  Монолайнер \u2014 только инвестиции",
            "S5  Лидер инноваций (крипта, 24/7)",
            "S6  Экосистема полного цикла",
            "S7  25,4% оборотов MOEX (рекорд)",
        ]),
        ("WEAKNESSES", RED, 5.0, 1.05, [
            "W1  5-е место по узнаваемости (21%)",
            "W2  Конфликт: fintech внутри \u2260 образ",
            "W3  Низкая эмоц. близость (Likeability)",
            "W4  Коммуникация массовая и корпорат.",
            "W5  Восприятие \u00ABсложно\u00BB и \u00ABдавят\u00BB",
            "W6  Доля в кошельке \u21934,5 п.п.",
            "W7  Приток новых клиентов \u219323%",
        ]),
        ("OPPORTUNITIES", BLUE, 0.5, 3.35, [
            "O1  Архетип Маг/Искатель \u2014 свободен!",
            "O2  Сегмент 2М+ растёт (+0,11 п.п.)",
            "O3  Запрос: Sales \u2192 Partner/Enabler",
            "O4  Спрос на пассивный доход (56\u201360%)",
            "O5  Горизонт 3+ года (47\u201351% ЦА)",
            "O6  Fintech = топ-сектор для инвесторов",
            "O7  Независимость от банка = уник. образ",
        ]),
        ("THREATS", AMBER, 5.0, 3.35, [
            "T1  Банки доминируют: 93,5% кл.",
            "T2  \u00ABСвязан с банком\u00BB = топ-фактор",
            "T3  Финам \u2014 конкурент, двойные счета",
            "T4  Недоверие к фин. советникам",
            "T5  Геополитическая неопределённость",
            "T6  АТОН усиливается в HNWI",
            "T7  Эмоц. барьеры: страх потерь",
        ]),
    ]
    for title, color, x, y, items in swot_data:
        # Header bar
        r = rect(s, x, y, 4.3, 0.35, color)
        txt(s, x, y + 0.02, 4.3, 0.3, title, FD, 13, WHITE, True, PP_ALIGN.CENTER)
        # Items
        for j, item in enumerate(items):
            txt(s, x + 0.15, y + 0.4 + j * 0.27, 4.0, 0.25, item, FT, 9, DGRAY)

    # ── SLIDE 5: Competitive Map (PNG — complex chart) ──────────────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_7.png")
    header(s, "01", "Конкурентная карта")
    img(s, FIG_DIR / "competitive_map.png", 0.5, 0.9, width=9, height=4.5)

    # ── SLIDE 6: Hypothesis Choice ──────────────────────────────────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_4.png")
    header(s, "02", "Выбор гипотезы", WHITE)
    txt(s, 1.2, 1.3, 7.5, 0.7, "Активатор благосостояния", FD, 40, WHITE, True)
    txt(s, 1.2, 2.1, 7.5, 0.35, "Архетип: Маг (65%) \u00D7 Искатель (35%)", FT, 16, WHITE)
    txt(s, 1.2, 2.6, 7.5, 0.8, "Суть: Мы превращаем ваши амбиции в работающий капитал, чтобы вы могли позволить себе всё, что для вас важно.", FT, 14, WHITE, lsp=1.4)
    for j, line in enumerate([
        "Почему эта гипотеза:",
        "\u2022 Архетипное пространство Маг/Искатель свободно на рынке РФ",
        "\u2022 Решает ключевой разрыв: рациональная сила есть, эмоциональной близости нет",
        "\u2022 Переводит бренд из \u00ABагрессивный продавец\u00BB в \u00ABтехнологичный партнёр\u00BB",
    ]):
        bold = j == 0
        txt(s, 1.2, 3.5 + j * 0.4, 7.5, 0.35, line, FT, 12, WHITE, bold)

    # ── SLIDE 7: Brand Platform (PNG — complex concentric) ──────────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_7.png")
    header(s, "02", "Платформа бренда")
    img(s, FIG_DIR / "brand_wheel.png", 1.5, 0.8, height=4.7)

    # ── SLIDE 8: Values + Big Idea ──────────────────────────────────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_7.png")
    header(s, "02", "Ценности бренда")
    txt(s, 0.8, 1.0, 5.5, 0.7, "Миссия: Превращаем финансовые амбиции в работающий капитал \u2014 чтобы деньги создавали жизнь, а не жизнь вращалась вокруг денег.", FT, 11, DGRAY, lsp=1.4)
    for i, (name, desc, color, role) in enumerate([
        ("ЯСНОСТЬ", "Сложность рынков \u2192 ясность решений", BLUE, "Фундамент"),
        ("ТРАНСФОРМАЦИЯ", "Потенциал \u2192 работающий капитал", VIOLET, "Двигатель"),
        ("МАСТЕРСТВО", "30 лет экспертизы \u2192 точность решений", GREEN, "Доказательство"),
        ("СВОБОДА", "Инвестиции = средство. Жизнь = цель", AMBER, "Обещание"),
    ]):
        y = 2.1 + i * 0.7
        rect(s, 0.8, y, 0.15, 0.45, color)  # color accent bar
        txt(s, 1.1, y, 2.0, 0.3, name, FD, 15, color, True)
        txt(s, 3.2, y + 0.05, 3.5, 0.3, desc, FT, 11, DGRAY)
        txt(s, 6.8, y + 0.05, 1.5, 0.3, role, FT, 9, LGRAY, align=PP_ALIGN.RIGHT)
    # Big Idea
    txt(s, 6.8, 0.35, 2.8, 0.3, "Big Idea:", FT, 11, DGRAY)
    txt(s, 6.8, 0.65, 2.8, 0.6, "Активируй\nсвоё благосостояние", FD, 20, GOLD, True, lsp=1.1)

    # ── SLIDE 9: Brand Architecture ─────────────────────────────────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_7.png")
    header(s, "02", "Бренд-архитектура")
    txt(s, 0.8, 1.0, 8.4, 0.3, "Endorsed Monolith: единая Big Idea \u00ABАктивируй\u00BB + объект = суббренд", FT, 12, DGRAY)

    # Parent
    r = rect(s, 2.5, 1.5, 5, 0.7, GOLD)
    txt(s, 2.5, 1.5, 5, 0.35, "БКС Мир Инвестиций", FD, 16, BLACK, True, PP_ALIGN.CENTER)
    txt(s, 2.5, 1.85, 5, 0.3, "\u00ABАктивируй благосостояние\u00BB", FT, 11, BLACK, align=PP_ALIGN.CENTER)

    for i, (name, slogan, vals) in enumerate([
        ("Платформа", "\u00ABАктивируй мастерство\u00BB", "Ясность + Трансформация"),
        ("Экспертиза", "\u00ABАктивируй стратегию\u00BB", "Мастерство + Трансформация"),
        ("Ультима", "\u00ABАктивируй свободу\u00BB", "Свобода + Трансформация"),
    ]):
        x = 0.8 + i * 3.1
        rect(s, x, 2.6, 2.8, 0.9, WHITE)
        txt(s, x, 2.65, 2.8, 0.3, name, FD, 15, BLACK, True, PP_ALIGN.CENTER)
        txt(s, x, 2.95, 2.8, 0.3, slogan, FT, 11, GOLD, align=PP_ALIGN.CENTER)
        txt(s, x, 3.25, 2.8, 0.25, vals, FT, 9, LGRAY, align=PP_ALIGN.CENTER)

    txt(s, 0.8, 3.8, 8.4, 0.35, "Медиа: БКС Экспресс \u00ABАктивируй знание\u00BB  \u00B7  БКС Профит \u00ABАктивируй стратегию\u00BB  \u00B7  Fintarget \u00ABАктивируй автопилот\u00BB", FT, 10, DGRAY)
    txt(s, 0.8, 4.3, 8.4, 0.35, "Трансформация \u2014 сквозная ценность через все суббренды (ДНК Мага)", FT, 11, VIOLET, True)

    # ── SLIDE 10: Comms Platform (NATIVE PPTX) ──────────────────────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_7.png")
    header(s, "03", "Коммуникационная платформа")

    # Big Idea bar
    r = rect(s, 2.5, 1.0, 5, 0.6, GOLD)
    txt(s, 2.5, 1.05, 5, 0.35, "АКТИВИРУЙ СВОЁ БЛАГОСОСТОЯНИЕ", FD, 14, BLACK, True, PP_ALIGN.CENTER)
    txt(s, 2.5, 1.35, 5, 0.2, "Big Idea", FT, 8, DGRAY, align=PP_ALIGN.CENTER)

    msgs = [
        ("Сообщение 1", "ЯСНОСТЬ", "Сложность рынков \u2192\nясность решений", "Сложный рынок.\nЯсные решения.", BLUE),
        ("Сообщение 2", "ТРАНСФОРМАЦИЯ", "Мы не храним капитал \u2014\nмы активируем его", "Капитал\nв действии.", VIOLET),
        ("Сообщение 3", "СВОБОДА", "30 лет экспертизы \u2014\nчтобы деньги работали\nна вашу жизнь", "Инвестируй в жизнь,\nа не в тревогу.", AMBER),
    ]
    for i, (label, value, body, slogan, color) in enumerate(msgs):
        x = 0.5 + i * 3.15
        # Card border
        shape = slide_rect = rect(s, x, 1.85, 2.9, 2.8, WHITE)
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)
        # Color bar
        rect(s, x + 0.05, 1.9, 2.8, 0.3, color)
        txt(s, x, 1.92, 2.9, 0.25, label, FD, 9, WHITE, True, PP_ALIGN.CENTER)
        # Value
        txt(s, x, 2.3, 2.9, 0.3, value, FD, 14, color, True, PP_ALIGN.CENTER)
        # Body
        txt(s, x + 0.15, 2.7, 2.6, 0.8, body, FT, 9, DGRAY, align=PP_ALIGN.CENTER, lsp=1.3)
        # Slogan
        txt(s, x + 0.15, 3.7, 2.6, 0.25, "Слоган:", FT, 7, LGRAY, align=PP_ALIGN.CENTER)
        txt(s, x + 0.15, 3.9, 2.6, 0.6, slogan, FD, 10, BLACK, True, PP_ALIGN.CENTER, lsp=1.2)

    # Funnel
    txt(s, 0.8, 4.85, 8.4, 0.25, "Awareness: Трансформация (наружка)  \u2192  Consideration: Ясность (контент)  \u2192  Conversion: Свобода (CTA)", FT, 9, DGRAY, align=PP_ALIGN.CENTER)

    # ── SLIDE 11: Slogans + Ad Mockups ──────────────────────────────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_7.png")
    header(s, "03", "Слоганы и макеты")

    for i, (num, slogan, ctx) in enumerate([
        ("1.", "Активируй своё благосостояние", "Основной таглайн"),
        ("2.", "Капитал в действии", "Имидж, наружка"),
        ("3.", "Сложный рынок. Ясные решения.", "Продукт, аналитика"),
        ("4.", "30 лет. Одна цель \u2014 ваше благосостояние.", "Heritage, PR"),
        ("5.", "Инвестируй в жизнь, а не в тревогу.", "Lifestyle, awareness"),
    ]):
        y = 1.0 + i * 0.35
        txt(s, 0.8, y, 0.3, 0.25, num, FT, 10, GOLD, True)
        txt(s, 1.2, y, 3.8, 0.25, slogan, FD, 12, BLACK, True)
        txt(s, 5.2, y, 2, 0.25, ctx, FT, 9, LGRAY)

    img(s, FIG_DIR / "ad_mockups.png", 0.3, 2.85, width=9.4, height=2.5)

    # ── SLIDE 12: Color System (PNG — swatches) ─────────────────────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_7.png")
    header(s, "04", "Цветовая система")
    img(s, FIG_DIR / "color_system.png", 0.3, 0.95, width=9.4, height=4.4)

    # ── SLIDE 13: Typography + Graphic Device (NATIVE + small PNG) ──────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_7.png")
    header(s, "04", "Типографика и графический приём")

    # Typography (left half — native)
    txt(s, 0.8, 1.1, 4.5, 0.25, "ТИПОГРАФИКА", FD, 11, LGRAY, True)
    specimens = [
        ("H1 Hero", FD, 24, "Активируй", BLACK, True),
        ("H2 Section", FD, 18, "Капитал в действии", BLACK, True),
        ("Body", FT, 12, "Мы превращаем амбиции в капитал.", DGRAY, False),
        ("Caption", FT, 9, "Рейтинг ruА-, ААА|ru.iv|. Данные 2026.", LGRAY, False),
        ("Data", FD, 16, "856 млрд \u20BD  \u00B7  25,4%", GOLD, True),
    ]
    y = 1.45
    for label, font, sz, text, color, bold in specimens:
        txt(s, 0.8, y, 0.8, 0.2, label, FT, 7, LGRAY)
        txt(s, 1.8, y, 3.5, 0.3, text, font, sz, color, bold)
        y += 0.6 if sz > 14 else 0.45

    # Graphic device (right half — PNG, intentionally dark)
    txt(s, 5.5, 1.1, 4, 0.25, "ГРАФИЧЕСКИЙ ПРИЁМ", FD, 11, LGRAY, True)
    txt(s, 5.5, 1.4, 4, 0.25, "Activation Pulse \u2014 импульс активации", FT, 10, DGRAY)
    img(s, FIG_DIR / "graphic_device.png", 5.3, 1.8, width=4.5, height=3.3)

    # ── SLIDE 14: Applied Identity (NATIVE layout + small PNGs) ─────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_7.png")
    header(s, "04", "Применение айдентики")
    txt(s, 0.8, 0.9, 8.4, 0.25, "Digital banner  \u00B7  Наружная реклама  \u00B7  Интерфейс приложения", FT, 11, DGRAY)
    img(s, FIG_DIR / "applied_identity.png", 0.3, 1.2, width=9.4, height=4.2)

    # ── SLIDE 15: Thank You ─────────────────────────────────────────────
    s = prs.slides.add_slide(B)
    add_bg(s, prs, "background_2.png")
    txt(s, 1.5, 1.2, 7, 1, "Активируй\nсвоё благосостояние", FD, 44, WHITE, True, PP_ALIGN.CENTER, lsp=1.15)
    txt(s, 1.5, 2.8, 7, 0.4, "БКС Мир Инвестиций  \u00B7  Стратегия репозиционирования", FT, 14, WHITE, align=PP_ALIGN.CENTER)
    txt(s, 1.5, 3.6, 7, 0.4, "DDVB  \u00B7  Direct Design Visual Branding", FD, 16, WHITE, align=PP_ALIGN.CENTER)
    txt(s, 1.5, 4.1, 7, 0.35, "ddvb.ru  \u00B7  +7 495 916 0123  \u00B7  Март 2026", FT, 11, WHITE, align=PP_ALIGN.CENTER)

    # ── Save ────────────────────────────────────────────────────────────
    out = OUT_DIR / "BCS_Repositioning_DDVB_2026_v4.pptx"
    prs.save(str(out))
    mb = out.stat().st_size / (1024 * 1024)
    print(f"Saved: {out}")
    print(f"Size: {mb:.1f} MB  |  Slides: {len(prs.slides)}  |  Limit: 15 MB")


if __name__ == "__main__":
    build()
